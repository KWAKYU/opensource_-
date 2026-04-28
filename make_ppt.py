from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 (PDF 디자인 기준) ──────────────────
ACCENT    = RGBColor(0xE8, 0x40, 0x55)  # 로즈 레드 (주 강조색)
ACCENT_LT = RGBColor(0xFF, 0xF0, 0xF3)  # 연핑크 (카드 배경)
DARK      = RGBColor(0x1A, 0x1A, 0x2E)  # 다크 텍스트
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x77, 0x77, 0x77)
GRAY_LT   = RGBColor(0xF5, 0xF5, 0xF5)  # 연회색 카드

RECT  = 1  # 직사각형
RRECT = 5  # 모서리 둥근 직사각형
OVAL  = 9  # 타원/원


def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height,
              fill_color=ACCENT, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 0.75)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                size=16, bold=False, color=DARK,
                align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  size=15, color=DARK, line_spacing_pt=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if line_spacing_pt:
            from pptx.util import Pt as _Pt
            from pptx.oxml.ns import qn
            import lxml.etree as etree
            pPr = p._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
            spcPts.set('val', str(int(line_spacing_pt * 100)))
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox


def slide_title(slide, title):
    """슬라이드 공통 제목: 왼쪽 세로 바 + 굵은 제목"""
    add_shape(slide, RECT, 0.42, 0.30, 0.07, 0.58, ACCENT)
    add_textbox(slide, title, 0.65, 0.26, 11.5, 0.75, size=30, bold=True, color=DARK)


def number_badge(slide, num_text, cx, cy, size_in=0.28):
    """빨간 원형 번호 배지"""
    add_shape(slide, OVAL, cx - size_in/2, cy - size_in/2,
              size_in, size_in, fill_color=ACCENT)
    add_textbox(slide, num_text,
                cx - size_in/2 - 0.02, cy - size_in/2 - 0.01,
                size_in + 0.04, size_in + 0.02,
                size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def card(slide, left, top, width, height,
         num=None, title="", subtitle="",
         bg=ACCENT_LT, border=ACCENT,
         title_size=15, sub_size=12):
    """번호 배지 + 제목 + 설명이 있는 카드"""
    add_shape(slide, RRECT, left, top, width, height,
              fill_color=bg, line_color=border, line_width=0.75)
    x_off = left + 0.12
    if num:
        number_badge(slide, num, left + 0.22, top + 0.22)
        x_off = left + 0.48
    if title:
        add_textbox(slide, title, x_off, top + 0.08,
                    width - x_off + left - 0.08, 0.35,
                    size=title_size, bold=True, color=DARK)
    if subtitle:
        add_textbox(slide, subtitle, left + 0.12, top + 0.42,
                    width - 0.24, height - 0.52,
                    size=sub_size, color=GRAY)


# ════════════════════════════════════════════════
# 슬라이드 1 — 타이틀 (배경 전체 ACCENT)
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, ACCENT)

# 좌상단 작은 흰 줄
add_shape(slide, RECT, 0.6, 0.38, 0.38, 0.07, WHITE)

# 우상단 정보
add_textbox(slide, "오픈소스 SW 응용 중간과제",
            8.5, 0.25, 4.5, 0.35, size=12, color=WHITE, align=PP_ALIGN.RIGHT)
add_textbox(slide, "곽유 / 2022106065",
            8.5, 0.58, 4.5, 0.35, size=12, color=WHITE, align=PP_ALIGN.RIGHT)
add_textbox(slide, "2026-04-26",
            8.5, 0.91, 4.5, 0.35, size=12, color=WHITE, align=PP_ALIGN.RIGHT)

# 메인 타이틀 (하단 왼쪽)
add_textbox(slide, "서울 코스\n추천기: AI\n멀티에이전트 토론 시스템",
            0.6, 2.5, 9.5, 3.5, size=52, bold=True, color=WHITE)
add_textbox(slide, "AI 에이전트들이 토론해서 만드는 나들이 코스",
            0.6, 6.1, 10, 0.55, size=18, color=WHITE)

# ════════════════════════════════════════════════
# 슬라이드 2 — 기존 서비스의 한계와 새로운 접근
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "기존 추천 서비스의 한계와 새로운 접근 방식")

# 왼쪽: 문제점 카드 3개
problems = [
    ("❶", "단순 큐레이션 의존",
     "카카오맵·데이트팝 등 기존 앱은 고정된\n추천 리스트와 큐레이션에 크게 의존함"),
    ("❷", "개인화 변수 고려 부족",
     "사용자의 예산, 동선, 테마를 동시에 정밀하게\n반영한 동적 경로 생성이 어려움"),
    ("❸", "단독 AI의 품질 한계",
     "AI 단독 생성물은 사실 검토와 신뢰성\n확보를 위한 다각도 검증 단계가 미흡함"),
]
y_pos = [1.20, 2.80, 4.40]
for (num, title, sub), y in zip(problems, y_pos):
    card(slide, 0.45, y, 5.8, 1.45, num=num, title=title, subtitle=sub)

# 오른쪽: 새로운 접근 박스
add_shape(slide, RRECT, 6.7, 1.10, 6.2, 5.80,
          fill_color=ACCENT_LT, line_color=ACCENT, line_width=1.0)
add_textbox(slide, "새로운 접근: 멀티에이전트 토론",
            6.95, 1.30, 5.7, 0.45, size=15, bold=True, color=ACCENT)

# 간단한 플로 다이어그램
add_shape(slide, RRECT, 7.0, 2.0, 2.2, 0.8, fill_color=ACCENT)
add_textbox(slide, "Agent A", 7.0, 2.07, 2.2, 0.5,
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_shape(slide, RRECT, 10.7, 2.0, 2.2, 0.8, fill_color=ACCENT)
add_textbox(slide, "Agent B", 10.7, 2.07, 2.2, 0.5,
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# 화살표 텍스트
add_textbox(slide, "제안  →", 9.25, 1.85, 1.4, 0.35,
            size=11, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(slide, "←  반박", 9.25, 2.45, 1.4, 0.35,
            size=11, color=DARK, align=PP_ALIGN.CENTER)

add_shape(slide, RRECT, 8.2, 3.2, 3.5, 0.8, fill_color=ACCENT)
add_textbox(slide, "최적의 서울 코스",
            8.2, 3.27, 3.5, 0.5, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "다각도 검증 완료",
            8.2, 3.68, 3.5, 0.35, size=10, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide, "가설: AI 간 토론을 통한 품질 고도화",
            6.95, 4.3, 5.7, 0.35, size=12, bold=True, color=DARK)
add_textbox(slide, "에이전트 간의 상호 검증으로 신뢰도 높은 코스 제안",
            6.95, 4.65, 5.7, 0.45, size=11, color=GRAY)

# ════════════════════════════════════════════════
# 슬라이드 3 — 멀티에이전트 토론 개요
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "멀티에이전트 토론: 집단 지성형 코스 설계")

features = [
    ("❶", "5명의 전문 AI 에이전트 협력 체계",
     "Planner, Scout, Budget, Experience, Verifier"),
    ("❷", "반복 토론을 통한 품질 고도화",
     "최소 2라운드 ~ 최대 5라운드 순환 구조"),
    ("❸", "엄격한 최종 합의 도달 조건",
     "만족도 8점 이상 및 예산 승인 필수"),
    ("❹", "비판적 검토 기반 오류 최소화",
     "단순 생성을 넘어선 검증 프로세스"),
]
y_pos = [1.15, 2.45, 3.75, 5.05]
for (num, title, sub), y in zip(features, y_pos):
    card(slide, 0.45, y, 5.8, 1.15, num=num, title=title, subtitle=sub)

# 오른쪽: 에이전트 원형 배치 (텍스트로 표현)
cx, cy = 10.0, 4.0
add_shape(slide, OVAL, cx - 1.0, cy - 1.0, 2.0, 2.0, fill_color=ACCENT)
add_textbox(slide, "MULTI-\nAGENT\nDEBATE",
            cx - 0.9, cy - 0.75, 1.8, 1.5,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

agents_around = [
    ("PLANNER",   cx,       1.15),
    ("SCOUT",     cx + 2.1, cy - 0.15),
    ("BUDGET",    cx + 1.3, cy + 1.8),
    ("EXPERIENCE",cx - 2.0, cy + 1.8),
    ("VERIFIER",  cx - 2.6, cy - 0.15),
]
for name, ax, ay in agents_around:
    add_textbox(slide, name, ax - 0.7, ay, 1.4, 0.38,
                size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_textbox(slide, "Round 2 ~ 5 Iteration",
            7.5, 6.55, 5.0, 0.38, size=12, bold=True, color=ACCENT)
add_textbox(slide, "Consensus Target: Score ≥ 8.0",
            7.5, 6.90, 5.0, 0.35, size=11, color=GRAY)

# ════════════════════════════════════════════════
# 슬라이드 4 — 전문화된 에이전트별 역할 분담
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "전문화된 에이전트별 역할 분담 체계")

agents = [
    ("❶", "Planner  (Claude Sonnet)",
     "전체적인 전략 수립 및 카테고리별 스케줄 설계"),
    ("❷", "Scout  (GPT-4o mini)",
     "Kakao API를 활용한 실제 장소 탐색 및 데이터 수집"),
    ("❸", "Budget  (Perplexity)",
     "실시간 웹 검색을 통한 실제 이용 가격 추정"),
    ("❹", "Experience  (Gemini 2.0 Flash)",
     "네이버 블로그 후기 분석 및 코스 비판/반박"),
    ("❺", "Verifier  (Claude Haiku)",
     "전체 토론 내용을 종합하여 최종 코스 확정"),
]

positions = [
    (0.45, 1.15), (6.95, 1.15),
    (0.45, 3.00), (6.95, 3.00),
    (0.45, 4.85),
]
for (num, title, sub), (lx, ly) in zip(agents, positions):
    card(slide, lx, ly, 5.95, 1.65, num=num, title=title, subtitle=sub,
         title_size=14, sub_size=12)

# ════════════════════════════════════════════════
# 슬라이드 5 — 토론 메커니즘: Experience → Scout 피드백
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "비판과 수용을 통한 실질적인 토론 메커니즘")

# 왼쪽: 플로 박스
add_shape(slide, RRECT, 0.45, 1.15, 5.8, 5.3,
          fill_color=GRAY_LT, line_color=GRAY, line_width=0.5)

add_shape(slide, RRECT, 0.7, 1.45, 2.3, 0.85, fill_color=ACCENT_LT, line_color=ACCENT)
add_textbox(slide, "Experience", 0.7, 1.54, 2.3, 0.38, size=13, bold=True, color=ACCENT,
            align=PP_ALIGN.CENTER)
add_textbox(slide, "반박(Objection) 제시", 0.7, 1.88, 2.3, 0.30, size=10, color=GRAY,
            align=PP_ALIGN.CENTER)

add_shape(slide, RRECT, 3.65, 1.45, 2.1, 0.85, fill_color=ACCENT_LT, line_color=ACCENT)
add_textbox(slide, "Scout", 3.65, 1.54, 2.1, 0.38, size=13, bold=True, color=ACCENT,
            align=PP_ALIGN.CENTER)
add_textbox(slide, "반박 내용 반영", 3.65, 1.88, 2.1, 0.30, size=10, color=GRAY,
            align=PP_ALIGN.CENTER)

# 화살표 텍스트
add_textbox(slide, "──  비판 내용 전달  ──→", 0.55, 1.58, 3.2, 0.32,
            size=9, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(slide, "←── 더 나은 장소로 교체 및 재탐색 ──",
            0.55, 2.48, 5.4, 0.32, size=9, color=GRAY, align=PP_ALIGN.CENTER)

add_textbox(slide, "실제 반영 사례:", 0.75, 2.95, 5.2, 0.35, size=12, bold=True, color=ACCENT)
add_shape(slide, RRECT, 0.7, 3.30, 5.3, 1.0,
          fill_color=WHITE, line_color=ACCENT, line_width=0.5)
add_textbox(slide, '"해당 카페는 블로그 후기가 너무 적습니다."',
            0.85, 3.38, 5.0, 0.40, size=12, bold=True, color=DARK)
add_textbox(slide, "→ 인기도 기반의 검증된 명소로 재탐색",
            1.1, 3.78, 4.8, 0.38, size=11, color=GRAY)

# 오른쪽: 특징 카드 2개 + 핵심 동력 박스
card(slide, 6.7, 1.15, 6.2, 1.95,
     num="❶", title="반박 중심의 운영 방식",
     subtitle="단순 독립 선택에서 벗어나 Experience의 비판을\n중심으로 한 상호 보완적 의사결정을 수행합니다.",
     title_size=14, sub_size=11)

card(slide, 6.7, 3.25, 6.2, 1.95,
     num="❷", title="에이전트 간 피드백 루프",
     subtitle="Experience의 비판이 다음 라운드 Scout에게\n직접 전달되어 코스 품질을 즉각 개선합니다.",
     title_size=14, sub_size=11)

add_shape(slide, RRECT, 6.7, 5.30, 6.2, 1.10,
          fill_color=ACCENT, line_color=None)
add_textbox(slide, "핵심 동력",
            6.95, 5.38, 5.8, 0.38, size=13, bold=True, color=WHITE)
add_textbox(slide, "에이전트 간 정교한 피드백 루프 형성이\n전체 시스템의 코스 품질 향상을 견인합니다.",
            6.95, 5.72, 5.8, 0.55, size=11, color=WHITE)

# ════════════════════════════════════════════════
# 슬라이드 6 — 데이터 기반의 장소 품질 고도화 전략
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "데이터 기반의 장소 품질 고도화 전략")

# 왼쪽: 비교 테이블
add_shape(slide, RRECT, 0.45, 1.15, 6.4, 5.3,
          fill_color=GRAY_LT, line_color=GRAY, line_width=0.5)

headers = ["구분", "단순 거리순", "품질 고도화(Proposed)"]
col_x   = [0.55, 1.55, 3.5]
col_w   = [0.95, 1.85, 3.0]
add_shape(slide, RECT, 0.45, 1.15, 6.4, 0.48, fill_color=GRAY_LT)
for h, x, w in zip(headers, col_x, col_w):
    bold = (h == "품질 고도화(Proposed)")
    color = ACCENT if bold else DARK
    add_textbox(slide, h, x, 1.20, w, 0.38, size=12, bold=bold, color=color)

rows = [
    ("추천 대상", "역내 편의시설 위주", "검증된 핫플레이스"),
    ("데이터 소스", "좌표값 기준", "API + 블로그 교차"),
    ("신뢰도", "낮음 (오분류 포함)", "매우 높음 (필터링)"),
]
for i, (label, before, after) in enumerate(rows):
    y = 1.75 + i * 0.82
    add_textbox(slide, label, col_x[0], y + 0.05, col_w[0], 0.55, size=12, color=DARK)
    add_textbox(slide, before, col_x[1], y + 0.05, col_w[1], 0.55, size=11, color=GRAY)
    color_after = ACCENT if i == 2 else DARK
    bold_after  = (i == 2)
    add_textbox(slide, after, col_x[2], y + 0.05, col_w[2], 0.55,
                size=11, bold=bold_after, color=color_after)
    if i < 2:
        add_shape(slide, RECT, 0.5, y + 0.72, 6.2, 0.02, fill_color=GRAY)

# 오른쪽: 핵심 해결책 3가지 + 전략 박스
add_textbox(slide, "단순 거리순 정렬의 문제점(지하철역 내 분식집 위주) 해결",
            6.95, 1.20, 6.0, 0.50, size=12, bold=True, color=DARK)

solutions_q = [
    ("❶", "Kakao Local API 정확도순 적용",
     "좌표+반경 기반 검색 시 sort=distance 활용으로\n지역 중심의 검색 품질 확보"),
    ("❷", "블로그 언급수 기반 인기도 정량화",
     "네이버 블로그 언급 1,000건 이상 = 검증된 핫플\n100건 미만 = 무명 장소 제외"),
    ("❸", "오분류 데이터 자동 필터링",
     "음식점 키워드 감지 → 포토스팟/쇼핑 오분류 자동\n수정 (_sanitize_course 후처리)"),
]
for i, (num, title, sub) in enumerate(solutions_q):
    card(slide, 6.95, 1.78 + i * 1.65, 6.0, 1.50,
         num=num, title=title, subtitle=sub, title_size=13, sub_size=10)

add_shape(slide, RRECT, 6.95, 6.60, 6.0, 0.80, fill_color=ACCENT_LT, line_color=ACCENT)
add_textbox(slide, "장소 신뢰도 확보 전략",
            7.15, 6.65, 5.6, 0.32, size=11, bold=True, color=DARK)
add_textbox(slide, "신뢰할 수 있는 데이터 소스 교차 검증을 통한 품질 고도화",
            7.15, 6.95, 5.6, 0.32, size=10, color=GRAY)

# ════════════════════════════════════════════════
# 슬라이드 7 — Docker Compose 기반 다중 서비스 배포
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Docker Compose 기반 다중 서비스 배포")

# 왼쪽: Compose 없을 때 vs 있을 때
add_textbox(slide, "Compose 없이 배포할 경우",
            0.45, 1.15, 5.9, 0.40, size=14, bold=True, color=DARK)
no_compose = [
    "① docker build -t app .",
    "② docker run streamlit (포트 8501)",
    "③ docker run api (포트 8000)",
    "④ docker run cli",
    "→ 옵션 4개 × 서비스 3개 = 12번 입력",
]
for i, line in enumerate(no_compose):
    color = ACCENT if "→" in line else GRAY
    bold  = "→" in line
    add_textbox(slide, line, 0.55, 1.60 + i * 0.42, 5.7, 0.38,
                size=12, bold=bold, color=color)

# 구분선
add_shape(slide, RECT, 0.45, 3.85, 5.9, 0.03, fill_color=GRAY)

add_textbox(slide, "docker-compose.yml 한 파일로 해결",
            0.45, 3.95, 5.9, 0.40, size=14, bold=True, color=ACCENT)
with_compose = [
    "docker-compose up --build",
    "→ 1개 이미지 빌드",
    "→ streamlit + api + cli 동시 실행",
    "→ 환경변수(.env)·포트·볼륨 자동 적용",
    "→ 코드 변경 시 --build 옵션 추가하면 됨",
]
for i, line in enumerate(with_compose):
    color = ACCENT if "→" in line else DARK
    bold  = "docker-compose" in line
    add_textbox(slide, line, 0.55, 4.40 + i * 0.42, 5.7, 0.38,
                size=12, bold=bold, color=color)

# 오른쪽: 3개 서비스 카드
add_textbox(slide, "1 Image — 3 Services",
            6.9, 1.15, 6.0, 0.42, size=16, bold=True, color=DARK)

services = [
    ("streamlit", "포트 8501", "에이전트 토론 과정 실시간 시각화 웹 UI"),
    ("api",       "포트 8000", "FastAPI REST 서버  /plan  /health 엔드포인트"),
    ("cli",       "터미널",    "터미널 직접 실행용 명령행 인터페이스"),
]
for i, (name, port, desc) in enumerate(services):
    y = 1.65 + i * 1.75
    add_shape(slide, RRECT, 6.9, y, 6.0, 1.55,
              fill_color=ACCENT_LT, line_color=ACCENT, line_width=0.75)
    add_textbox(slide, name, 7.1, y + 0.12, 2.2, 0.42, size=16, bold=True, color=ACCENT)
    add_textbox(slide, port, 9.5, y + 0.16, 1.5, 0.35, size=12, color=GRAY)
    add_textbox(slide, desc, 7.1, y + 0.58, 5.6, 0.75, size=12, color=DARK)

add_shape(slide, RRECT, 6.9, 6.82, 6.0, 0.55, fill_color=ACCENT)
add_textbox(slide, "변경 없을 시: up  /  코드 변경 후: up --build",
            6.9, 6.88, 6.0, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════
# 슬라이드 8 — 반복 개발 과정 (9회 반복)
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "반복 개발을 통한 시스템 진화  (총 9회 반복)")

iters = [
    ("1", "뼈대 구축",   "에이전트 5개 설계\n_parse_json 방어 파싱"),
    ("2", "API 연동",    "Kakao Local API\nlazy init 패턴 적용"),
    ("3", "인프라",      "FastAPI + Docker\ndocker-compose 구성"),
    ("4", "UI + 가격",   "Streamlit UI\nNaver Blog API 연동"),
    ("5", "인기도",      "blog_count 지표\n무명 장소 필터링"),
    ("6", "오분류 수정", "FOOD_KEYWORDS 필터\n주소 기반 지하철 제거"),
    ("7", "교정 로직",   "_sanitize_course\n카테고리 중복 제거"),
    ("8", "Rate Limit",  "Gemini 재시도 로직\n5s/10s/15s 백오프"),
    ("9", "토론 강화",   "previous_feedback\n실질적 반복 토론"),
]

col_w = 1.33
col_h = 5.5
for i, (num, title, desc) in enumerate(iters):
    x = 0.4 + i * (col_w + 0.04)
    # 배경 카드
    add_shape(slide, RRECT, x, 1.15, col_w, col_h,
              fill_color=ACCENT_LT, line_color=ACCENT, line_width=0.5)
    # 상단 색 바
    add_shape(slide, RRECT, x, 1.15, col_w, 0.55, fill_color=ACCENT)
    add_textbox(slide, f"반복 {num}", x + 0.05, 1.18, col_w - 0.1, 0.38,
                size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, title, x + 0.05, 1.82, col_w - 0.1, 0.42,
                size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, desc, x + 0.08, 2.35, col_w - 0.16, 4.0,
                size=10, color=GRAY)

# ════════════════════════════════════════════════
# 슬라이드 9 — 결론 및 핵심 성과
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "결론 및 핵심 성과")

# 왼쪽 성과 카드들
achievements = [
    ("❶", "다각도 검증 기반 품질 향상",
     "단독 AI 대비 오분류·무명 장소 추천 대폭 감소\n반복 토론으로 점수 8점 이상 코스 안정적 달성"),
    ("❷", "실시간 데이터 기반 신뢰도 확보",
     "Kakao + Naver Blog 이중 소스 교차 검증\n인기도 정량화로 '검증된 핫플' 우선 추천"),
    ("❸", "확장 가능한 멀티에이전트 아키텍처",
     "에이전트 교체·추가가 독립적으로 가능\nDocker Compose로 손쉬운 배포 환경 제공"),
]
for i, (num, title, sub) in enumerate(achievements):
    card(slide, 0.45, 1.15 + i * 2.0, 6.0, 1.80,
         num=num, title=title, subtitle=sub, title_size=14, sub_size=11)

# 오른쪽: 기술 스택 + 가산점
add_textbox(slide, "기술 스택 & 가산점",
            6.8, 1.15, 6.1, 0.42, size=15, bold=True, color=DARK)

stack_items = [
    ("OpenRouter",            "Claude / GPT-4o mini / Gemini / Perplexity"),
    ("Kakao Local API",       "국내 장소 탐색  (좌표·반경 기반)"),
    ("Naver Blog Search API", "인기도 정량화 및 분위기 후기 수집"),
    ("Streamlit",             "토론 과정 실시간 시각화 웹 UI"),
    ("FastAPI + Docker",      "REST API + 컨테이너화  (+5+5 가산점)"),
]
for i, (name, desc) in enumerate(stack_items):
    y = 1.65 + i * 0.68
    add_textbox(slide, f"• {name}", 6.85, y, 2.5, 0.38, size=12, bold=True, color=ACCENT)
    add_textbox(slide, desc, 9.45, y, 3.5, 0.38, size=11, color=GRAY)

# 배운 점 박스
add_shape(slide, RRECT, 6.8, 5.15, 6.1, 1.5,
          fill_color=ACCENT_LT, line_color=ACCENT, line_width=0.75)
add_textbox(slide, "핵심 교훈",
            7.0, 5.22, 5.8, 0.38, size=13, bold=True, color=ACCENT)
learnings = [
    "에이전트 역할 분리가 명확할수록 토론 품질 향상",
    "LLM 신뢰 + Python 후처리 이중 안전망으로 오류 최소화",
    "피드백 루프 설계가 시스템 전체 성능을 결정한다",
]
for i, line in enumerate(learnings):
    add_textbox(slide, f"• {line}", 7.0, 5.62 + i * 0.33, 5.8, 0.32,
                size=11, color=DARK)

# ════════════════════════════════════════════════
out = "/Users/ukwak/claude code/중간과제/서울코스추천기_발표.pptx"
prs.save(out)
print(f"저장 완료: {out}")
